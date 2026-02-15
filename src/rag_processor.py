import os
import PyPDF2
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import requests
import json
from typing import List, Tuple
import re
import base64
import ast

# For different file types
try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

class RAGProcessor:
    def __init__(self, embedding_model='all-MiniLM-L6-v2', data_dir='data'):
        self.embedding_model = SentenceTransformer(embedding_model)
        self.index = None
        self.documents = []
        self.document_paths = []  # Store full paths
        self.chunks = []
        self.chunk_metadata = []  # Store metadata for each chunk
        self.chunk_size = 1000
        self.overlap = 200
        self.data_dir = data_dir
        self.visual_features = {}  # Track which files have visual features
        os.makedirs(self.data_dir, exist_ok=True)
        
        # Try to load existing data
        self.load_data()

    def save_data(self):
        """Save the FAISS index and document data to disk."""
        try:
            if self.index is not None:
                # Save FAISS index
                faiss.write_index(self.index, os.path.join(self.data_dir, 'faiss_index.idx'))
                
                # Save document data
                data = {
                    'documents': self.documents,
                    'document_paths': self.document_paths,
                    'chunks': self.chunks,
                    'chunk_metadata': self.chunk_metadata,
                    'visual_features': self.visual_features
                }
                with open(os.path.join(self.data_dir, 'documents.json'), 'w') as f:
                    json.dump(data, f, indent=2)
        except Exception as e:
            print(f"Error saving data: {e}")

    def save_uploaded_files_list(self, uploaded_files_data):
        """Save the list of uploaded files and their metadata, avoiding duplicates."""
        try:
            # Load existing files
            existing_files = self.load_uploaded_files_list()
            
            # Create a set of existing file identifiers (name + size + lastModified)
            existing_ids = set()
            for file_data in existing_files:
                file_id = f"{file_data['name']}_{file_data['size']}_{file_data['lastModified']}"
                existing_ids.add(file_id)
            
            # Filter out duplicates
            new_files = []
            duplicates = 0
            
            for file_data in uploaded_files_data:
                file_id = f"{file_data['name']}_{file_data['size']}_{file_data['lastModified']}"
                if file_id not in existing_ids:
                    new_files.append(file_data)
                else:
                    duplicates += 1
            
            # Combine existing and new files
            all_files = existing_files + new_files
            
            # Save the combined list
            with open(os.path.join(self.data_dir, 'uploaded_files.json'), 'w') as f:
                json.dump(all_files, f, indent=2)
            
            return {
                'saved': len(new_files),
                'duplicates': duplicates,
                'total': len(all_files)
            }
            
        except Exception as e:
            print(f"Error saving uploaded files list: {e}")
            return {
                'saved': 0,
                'duplicates': 0,
                'total': len(existing_files) if 'existing_files' in locals() else 0,
                'error': str(e)
            }

    def load_uploaded_files_list(self):
        """Load the list of uploaded files and their metadata."""
        try:
            file_path = os.path.join(self.data_dir, 'uploaded_files.json')
            if os.path.exists(file_path):
                with open(file_path, 'r') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading uploaded files list: {e}")
        return []

    def load_data(self):
        """Load the FAISS index and document data from disk."""
        try:
            index_path = os.path.join(self.data_dir, 'faiss_index.idx')
            data_path = os.path.join(self.data_dir, 'documents.json')
            
            if os.path.exists(index_path) and os.path.exists(data_path):
                # Load FAISS index
                self.index = faiss.read_index(index_path)
                
                # Load document data
                with open(data_path, 'r') as f:
                    data = json.load(f)
                    self.documents = data.get('documents', [])
                    self.document_paths = data.get('document_paths', [])
                    self.chunks = data.get('chunks', [])
                    self.chunk_metadata = data.get('chunk_metadata', [])
                    self.visual_features = data.get('visual_features', {})
                    
                    # Ensure chunk_metadata has the same length as chunks (for backward compatibility)
                    if len(self.chunk_metadata) != len(self.chunks):
                        print(f"Warning: chunk_metadata length ({len(self.chunk_metadata)}) doesn't match chunks length ({len(self.chunks)}), initializing empty metadata")
                        self.chunk_metadata = [{'start_line': None, 'end_line': None, 'context': ''} for _ in self.chunks]
                    
                print(f"Loaded {len(self.documents)} documents from disk")
        except Exception as e:
            print(f"Error loading data: {e}")
            # Initialize empty if loading fails
            self.index = None
            self.documents = []
            self.document_paths = []
            self.chunks = []
            self.chunk_metadata = []

    def get_file_page_count(self, file_path: str) -> int:
        """Get the number of pages in a file (1 for non-PDF)."""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    return len(pdf_reader.pages)
            else:
                return 1  # For other files, consider as 1 page
        except Exception as e:
            print(f"Error getting page count for {file_path}: {e}")
            return 0

    def get_vlm_description(self, image_path: str, vlm_api_url: str, vlm_model: str) -> str:
        """Get description of an image using VLM API."""
        try:
            # Read and encode image to base64
            with open(image_path, 'rb') as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Determine MIME type
            ext = os.path.splitext(image_path)[1].lower()
            mime_type = f"image/{ext[1:]}"  # Remove the dot
            
            # Prepare the API request
            payload = {
                "model": vlm_model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this image in detail, including any text, objects, scenes, and context you can observe."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_data}"
                                }
                            }
                        ]
                    }
                ],
                "temperature": 0.7,
                "max_tokens": 500
            }
            
            response = requests.post(vlm_api_url, json=payload, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            description = result['choices'][0]['message']['content'].strip()
            return f"VLM Description: {description}"
            
        except Exception as e:
            print(f"Error getting VLM description for {image_path}: {e}")
            return ""

    def extract_docstrings(self, code: str) -> str:
        """Extract docstrings from Python code."""
        try:
            tree = ast.parse(code)
            docstrings = []
            for node in ast.walk(tree):
                if isinstance(node, (ast.FunctionDef, ast.ClassDef, ast.Module)) and ast.get_docstring(node):
                    docstrings.append(ast.get_docstring(node))
            return "\n".join(docstrings)
        except SyntaxError:
            return ""

    def extract_text_from_file(self, file_path: str, enable_vlm: bool = False, vlm_api_url: str = None, vlm_model: str = None) -> tuple[str, bool]:
        """Extract text from a file (PDF, DOCX, PPTX, XLSX, CSV, images, code files, text files). Returns (text, has_visual_features)."""
        text = ""
        has_visual_features = False
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
            elif ext == '.docx' and Document:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif ext == '.pptx' and Presentation:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext == '.xlsx' and openpyxl:
                wb = openpyxl.load_workbook(file_path)
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        text += " ".join(str(cell) for cell in row if cell) + "\n"
            elif ext == '.csv' and pd:
                df = pd.read_csv(file_path)
                text = df.to_string()
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                if enable_vlm and vlm_api_url and vlm_model:
                    # Use VLM for visual description
                    text = self.get_vlm_description(file_path, vlm_api_url, vlm_model)
                    has_visual_features = True
                elif pytesseract and Image:
                    # Fallback to OCR
                    img = Image.open(file_path)
                    text = pytesseract.image_to_string(img)
                else:
                    print(f"Neither VLM nor OCR available for image: {ext}")
            elif ext in ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.cs', '.php', '.rb', '.go', '.txt', '.md', '.rst', '.json', '.yaml', '.yml', '.xml', '.html', '.css', '.scss', '.sql']:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    if ext == '.py':
                        docstrings = self.extract_docstrings(text)
                        if docstrings:
                            text += "\n\nDocstrings:\n" + docstrings
                except Exception as e:
                    print(f"Error reading text file {file_path}: {e}")
                    text = ""
                has_visual_features = False
            else:
                print(f"Unsupported or missing library for file type: {ext}")
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
        return text, has_visual_features

    def chunk_text(self, text: str, file_path: str = None) -> List[dict]:
        """Split text into overlapping chunks with metadata."""
        chunks = []
        lines = text.split('\n')
        current_line = 1
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            chunk = text[start:end]
            
            # Calculate line numbers for this chunk
            chunk_start_line = current_line
            chunk_lines = chunk.split('\n')
            chunk_end_line = current_line + len(chunk_lines) - 1
            
            # Try to identify function/class context for code files
            context = self._get_code_context(text, chunk_start_line, chunk_end_line, file_path)
            
            chunks.append({
                'text': chunk,
                'start_line': chunk_start_line,
                'end_line': chunk_end_line,
                'context': context
            })
            
            # Update position and line counter for next chunk
            start = end - self.overlap
            # Count how many lines are in the overlap
            if start > 0:
                overlap_text = text[max(0, start):end]
                overlap_lines = overlap_text.split('\n')
                current_line = chunk_end_line - len(overlap_lines) + 1
            else:
                current_line = chunk_end_line + 1
        
        return chunks

    def _get_code_context(self, text: str, start_line: int, end_line: int, file_path: str = None) -> str:
        """Get code context (function/class name) for a given line range."""
        if not file_path:
            return ""
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.py':
            return self._get_python_context(text, start_line, end_line)
        else:
            # For other languages, try basic pattern matching
            return self._get_generic_code_context(text, start_line, end_line, ext)
    
    def _get_python_context(self, text: str, start_line: int, end_line: int) -> str:
        """Get Python function/class context using AST."""
        try:
            tree = ast.parse(text)
            context_stack = []
            
            for node in ast.walk(tree):
                if isinstance(node, (ast.FunctionDef, ast.ClassDef, ast.Module)):
                    # Check if this node spans the chunk's line range
                    if hasattr(node, 'lineno') and node.lineno <= end_line:
                        if hasattr(node, 'end_lineno') and node.end_lineno >= start_line:
                            if isinstance(node, ast.FunctionDef):
                                context_stack.append(f"def {node.name}")
                            elif isinstance(node, ast.ClassDef):
                                context_stack.append(f"class {node.name}")
                            elif isinstance(node, ast.Module) and node.name:
                                context_stack.append(f"module {node.name}")
            
            # Return the most specific context (innermost)
            return " > ".join(context_stack[-2:]) if context_stack else ""
        except:
            return ""
    
    def _get_generic_code_context(self, text: str, start_line: int, end_line: int, ext: str) -> str:
        """Get basic code context for non-Python files using regex patterns."""
        lines = text.split('\n')
        context = []
        
        # Look for function/class definitions in the chunk area
        for i in range(max(0, start_line - 10), min(len(lines), end_line + 10)):
            line = lines[i].strip()
            
            # JavaScript/TypeScript
            if ext in ['.js', '.ts', '.jsx', '.tsx']:
                if line.startswith(('function ', 'const ', 'let ', 'var ')) and ('=' in line or '(' in line):
                    if 'function' in line:
                        func_match = re.search(r'function\s+(\w+)', line)
                        if func_match:
                            context.append(f"function {func_match.group(1)}")
                    elif '=' in line and ('=>' in line or 'function' in line):
                        var_match = re.search(r'(?:const|let|var)\s+(\w+)', line)
                        if var_match:
                            context.append(f"function {var_match.group(1)}")
                elif line.startswith('class '):
                    class_match = re.search(r'class\s+(\w+)', line)
                    if class_match:
                        context.append(f"class {class_match.group(1)}")
            
            # Java/C#/C++
            elif ext in ['.java', '.cs', '.cpp', '.c', '.h']:
                if 'public ' in line or 'private ' in line or 'protected ' in line or 'static ' in line:
                    if '(' in line and ')' in line:  # Function
                        func_match = re.search(r'(?:public|private|protected|static)?\s*\w+\s+(\w+)\s*\(', line)
                        if func_match:
                            context.append(f"function {func_match.group(1)}")
                    elif 'class ' in line:
                        class_match = re.search(r'class\s+(\w+)', line)
                        if class_match:
                            context.append(f"class {class_match.group(1)}")
            
            # PHP
            elif ext == '.php':
                if line.startswith('function '):
                    func_match = re.search(r'function\s+(\w+)', line)
                    if func_match:
                        context.append(f"function {func_match.group(1)}")
                elif line.startswith('class '):
                    class_match = re.search(r'class\s+(\w+)', line)
                    if class_match:
                        context.append(f"class {class_match.group(1)}")
            
            # Ruby
            elif ext == '.rb':
                if line.startswith('def '):
                    func_match = re.search(r'def\s+(\w+)', line)
                    if func_match:
                        context.append(f"def {func_match.group(1)}")
                elif line.startswith('class '):
                    class_match = re.search(r'class\s+(\w+)', line)
                    if class_match:
                        context.append(f"class {class_match.group(1)}")
            
            # Go
            elif ext == '.go':
                if line.startswith('func '):
                    func_match = re.search(r'func\s+(?:\([^)]*\)\s*)?(\w+)', line)
                    if func_match:
                        context.append(f"func {func_match.group(1)}")
                elif line.startswith('type ') and 'struct' in line:
                    type_match = re.search(r'type\s+(\w+)\s+struct', line)
                    if type_match:
                        context.append(f"type {type_match.group(1)}")
        
        return " > ".join(context[-2:]) if context else ""

    def process_files(self, file_paths: List[str], enable_vlm: bool = False, vlm_api_url: str = None, vlm_model: str = None) -> bool:
        """Process multiple files: extract text, create chunks, and build vector index."""
        try:
            all_chunks = []
            all_chunk_metadata = []
            all_documents = []
            all_paths = []

            for file_path in file_paths:
                text, has_visual = self.extract_text_from_file(file_path, enable_vlm, vlm_api_url, vlm_model)
                if text.strip():
                    chunks = self.chunk_text(text, file_path)
                    all_chunks.extend([chunk['text'] for chunk in chunks])
                    all_chunk_metadata.extend([{
                        'start_line': chunk['start_line'],
                        'end_line': chunk['end_line'],
                        'context': chunk['context']
                    } for chunk in chunks])
                    all_documents.extend([os.path.basename(file_path)] * len(chunks))
                    all_paths.extend([file_path] * len(chunks))
                    
                    # Track visual features for this file
                    if has_visual:
                        self.visual_features[file_path] = True

            if not all_chunks:
                return False

            # Create embeddings
            embeddings = self.embedding_model.encode(all_chunks, show_progress_bar=True)

            # Build FAISS index
            dimension = embeddings.shape[1]
            self.index = faiss.IndexFlatL2(dimension)
            self.index.add(np.array(embeddings).astype('float32'))

            self.chunks = all_chunks
            self.chunk_metadata = all_chunk_metadata
            self.documents = all_documents
            self.document_paths = all_paths

            # Save data to disk
            self.save_data()

            return True
        except Exception as e:
            print(f"Error processing files: {e}")
            return False

    def remove_file(self, file_path: str) -> bool:
        """Remove a file and its embeddings from the index."""
        try:
            file_path = file_path.replace('/', os.sep).replace('\\', os.sep)
            
            # Check if file exists in the index
            if file_path in self.document_paths:
                # Find all indices for this file
                indices_to_remove = []
                for i, path in enumerate(self.document_paths):
                    if path == file_path:
                        indices_to_remove.append(i)
                
                if indices_to_remove:
                    # Remove from all lists (in reverse order to maintain indices)
                    for i in sorted(indices_to_remove, reverse=True):
                        del self.chunks[i]
                        del self.chunk_metadata[i]
                        del self.documents[i]
                        del self.document_paths[i]
                    
                    # Rebuild FAISS index if we have remaining documents
                    if self.chunks:
                        embeddings = self.embedding_model.encode(self.chunks, show_progress_bar=True)
                        dimension = embeddings.shape[1]
                        self.index = faiss.IndexFlatL2(dimension)
                        self.index.add(np.array(embeddings).astype('float32'))
                    else:
                        self.index = None
                    
                    # Remove from visual features if present
                    if file_path in self.visual_features:
                        del self.visual_features[file_path]
                    
                    # Save updated data
                    self.save_data()
            
            # Always return True since file deletion from disk is handled by the caller
            return True
        except Exception as e:
            print(f"Error removing file {file_path}: {e}")
            return False

    def has_visual_features(self, file_path: str) -> bool:
        """Check if a file has visual features (was processed with VLM)."""
        return file_path in self.visual_features

    def retrieve_relevant_chunks(self, query: str, top_k: int = 5, selected_files: List[str] = None) -> List[Tuple[str, str, str, float, dict]]:
        """Retrieve the most relevant chunks for a query, optionally filtered by selected files."""
        if self.index is None or not self.chunks:
            return []

        # Normalize selected_files to use OS path separator
        if selected_files:
            selected_files = [f.replace('/', os.sep).replace('\\', os.sep) for f in selected_files]

        query_embedding = self.embedding_model.encode([query])
        distances, indices = self.index.search(np.array(query_embedding).astype('float32'), top_k)

        results = []
        for i, idx in enumerate(indices[0]):
            if idx < len(self.chunks):
                chunk = self.chunks[idx]
                doc = self.documents[idx]
                path = self.document_paths[idx]
                metadata = self.chunk_metadata[idx] if idx < len(self.chunk_metadata) else {}
                
                # If selected_files is provided, only include chunks from selected files
                if selected_files and path not in selected_files:
                    continue
                    
                # Convert distance to similarity score (FAISS returns L2 distance, convert to cosine similarity)
                similarity_percent = float(f"{100.0 / (1.0 + distances[0][i]):.2f}")  # Simple conversion for display
                results.append((chunk, doc, path, similarity_percent, metadata))

        return results

    def chat(self, message: str, api_url: str, model: str, selected_files: List[str] = None, chat_history: List[dict] = None) -> dict:
        """Chat with the RAG system using the provided LLM API, optionally filtered by selected files."""
        # Retrieve relevant context
        relevant_chunks = self.retrieve_relevant_chunks(message, selected_files=selected_files)

        if not relevant_chunks:
            return {
                "answer": "No relevant information found in the selected files.",
                "sources": []
            }

        # Build context from relevant chunks (only for existing files)
        context_parts = []
        for chunk, doc, path, score, metadata in relevant_chunks:
            if os.path.exists(path):
                context_part = f"From {doc}: {chunk}"
                if metadata.get('context'):
                    context_part += f" (Context: {metadata['context']})"
                context_parts.append(context_part)
        context = "\n\n".join(context_parts)

        # Create unique sources list with additional metadata
        unique_sources = []
        seen_files = set()
        for chunk, doc, path, score, metadata in relevant_chunks:
            if os.path.exists(path) and doc not in seen_files:
                # Get file size and page count
                file_size = os.path.getsize(path) if os.path.exists(path) else 0
                page_count = self.get_file_page_count(path)
                
                source_info = {
                    "name": doc, 
                    "path": path, 
                    "file_size": file_size,
                    "page_count": page_count,
                    "similarity_percent": float(score)
                }
                
                # Add line number and context information if available
                if metadata.get('start_line'):
                    source_info["start_line"] = metadata['start_line']
                if metadata.get('end_line'):
                    source_info["end_line"] = metadata['end_line']
                if metadata.get('context'):
                    source_info["context"] = metadata['context']
                
                unique_sources.append(source_info)
                seen_files.add(doc)

        # Sort sources by similarity percent (highest first)
        unique_sources.sort(key=lambda x: x["similarity_percent"], reverse=True)

        # Prepare the prompt
        history_text = ""
        if chat_history:
            # Limit to last 10 messages to avoid token limits
            recent_history = chat_history[-10:]
            history_lines = []
            for msg in recent_history:
                sender = msg.get('sender', 'Unknown')
                text = msg.get('text', '')
                history_lines.append(f"{sender}: {text}")
            history_text = "\n".join(history_lines) + "\n\n"

        prompt = f"""You are a helpful assistant that answers questions based on the provided context from documents.

Context:
{context}

{history_text}Current question: {message}

Answer the question based on the provided context and conversation history. If the context doesn't contain enough information to answer the question, say so.

When referencing information from specific documents, use numbered citations like [1], [2], etc. corresponding to the sources provided."""

        # Call the LLM API
        try:
            payload = {
                "model": model,
                "messages": [
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7,
                "max_tokens": 1000
            }

            response = requests.post(api_url, json=payload, timeout=30)
            response.raise_for_status()

            result = response.json()
            answer = result['choices'][0]['message']['content'].strip()

            # Add href to sources
            for source in unique_sources:
                base_href = f"/uploads/{os.path.basename(source['path'])}"
                # Add line number anchor for code/text files
                if source.get('start_line'):
                    source["href"] = f"{base_href}#L{source['start_line']}"
                else:
                    source["href"] = base_href

            # Replace filename mentions with links
            for i, source in enumerate(unique_sources):
                number = str(i + 1)
                href = source["href"]
                filename = source["name"]
                answer = re.sub(r'\b' + re.escape(filename) + r'\b', f'<a href="{href}" target="_blank">[{number}]</a>', answer)

            # Replace citations with links
            for i, source in enumerate(unique_sources):
                number = str(i + 1)
                href = source["href"]
                answer = re.sub(r'\[' + re.escape(number) + r'\]', f'<a href="{href}" target="_blank">[{number}]</a>', answer)

            return {
                "answer": answer,
                "sources": unique_sources
            }

        except requests.exceptions.RequestException as e:
            return {
                "answer": f"Error calling LLM API: {str(e)}",
                "sources": []
            }
        except (KeyError, IndexError) as e:
            return {
                "answer": f"Error parsing LLM response: {str(e)}",
                "sources": []
            }